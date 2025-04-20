import os
import json
import logging
from huggingface_hub import InferenceClient
from django.conf import settings
from django.utils import timezone
import PyPDF2
import docx
import requests
from io import BytesIO
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not installed. PowerPoint files will not be processed.")

logger = logging.getLogger(__name__)

def get_hf_inference_client():
    """Get configured Hugging Face inference client"""
    api_key = settings.HUGGINGFACE_API_KEY
    
    # Log the API key (mask most of it for security)
    masked_key = api_key[:4] + "*****" + api_key[-4:] if len(api_key) > 8 else "****"
    logger.info(f"Initializing HF client with API key: {masked_key}")
    
    # Create and return the client with default provider
    # Using 'together' provider as primary for better reasoning capabilities
    client = InferenceClient(
        provider="together",
        api_key=api_key,
    )
    logger.info(f"HF client initialized: {client}")
    return client

def get_hf_fallback_client():
    """Get fallback HF client if primary fails"""
    api_key = settings.HUGGINGFACE_API_KEY
    logger.info("Initializing fallback HF client")
    
    # Using HF inference as fallback provider
    client = InferenceClient(
        provider="hf-inference",
        api_key=api_key,
    )
    logger.info(f"Fallback HF client initialized: {client}")
    return client

def generate_course_structure(prompt, context_documents=None):
    """
    Generate course title and modules based on the initial prompt
    """
    client = get_hf_inference_client()
    fallback_client = get_hf_fallback_client()
    
    # Create a context based on documents if provided
    context = ""
    if context_documents:
        context = "Based on the following context documents:\n\n"
        for doc in context_documents:
            context += f"{doc.content_text}\n\n"
    
    system_prompt = """You are an expert curriculum designer for cybersecurity courses. 
    Your task is to create a clear course structure with a title and 3-5 modules.
    Respond in JSON format with the following structure:
    {
        "course_title": "Title of the Course",
        "course_description": "A detailed description of the course",
        "course_type": "awareness" or "advanced",  
        "modules": [
            {
                "title": "Module 1 Title",
                "description": "Description of Module 1"
            },
            ...
        ]
    }
    The course_type should be "awareness" for general security awareness courses aimed at non-technical staff, 
    and "advanced" for technical security courses aimed at security professionals.
    """
    
    full_prompt = f"{system_prompt}\n\n{context}User request: {prompt}"
    
    try:
        # Try with primary model (Together AI with DeepSeek-R1)
        logger.info("Sending request to primary HF model (DeepSeek-R1)")
        logger.debug(f"Request model: deepseek-ai/DeepSeek-R1")
        logger.debug(f"Request provider: together")
        
        completion = client.chat.completions.create(
            model="deepseek-ai/DeepSeek-R1",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"{context}User request: {prompt}"}
            ],
            max_tokens=1000
        )
        response_text = completion.choices[0].message.content
        logger.info(f"Received response from primary model, length: {len(response_text)} characters")
        logger.debug(f"Response preview: {response_text[:100]}...")
    except Exception as e:
        logger.warning(f"Primary model failed, using fallback: {e}")
        # Fallback to HF inference model
        try:
            logger.info("Sending request to fallback HF model (Qwen)")
            logger.debug(f"Request model: Qwen/Qwen2-72B-Instruct")
            logger.debug(f"Request provider: hf-inference")
            
            completion = fallback_client.chat.completions.create(
                model="Qwen/Qwen2-72B-Instruct",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"{context}User request: {prompt}"}
                ],
                max_tokens=1000
            )
            response_text = completion.choices[0].message.content
            logger.info(f"Received response from fallback model, length: {len(response_text)} characters")
            logger.debug(f"Response preview: {response_text[:100]}...")
        except Exception as e2:
            logger.error(f"Both models failed: {e2}")
            return {"error": "Failed to generate course structure due to API issues."}
    
    # Try to extract JSON from the response
    try:
        # First attempt - direct parsing
        result = json.loads(response_text)
    except json.JSONDecodeError:
        # Second attempt - search for JSON in the text
        try:
            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            if json_start >= 0 and json_end > json_start:
                json_text = response_text[json_start:json_end]
                result = json.loads(json_text)
            else:
                raise ValueError("No valid JSON found in the response")
        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"Failed to parse JSON from response: {e}")
            return {"error": "Failed to parse the generated course structure."}
    
    return result

def generate_chapters(course_title, module_title, module_description, context_documents=None):
    """
    Generate chapters for a module based on the module title and description
    """
    client = get_hf_inference_client()
    fallback_client = get_hf_fallback_client()
    
    # Create a context based on documents if provided
    context = ""
    if context_documents:
        context = "Based on the following context documents:\n\n"
        for doc in context_documents:
            context += f"{doc.content_text}\n\n"
    
    system_prompt = f"""You are an expert curriculum designer for a cybersecurity course titled "{course_title}". 
    I need you to create chapters for the module titled "{module_title}" with description: "{module_description}".

    Create exactly 5 chapters for this module:
    - Chapters 1-4 should be teaching content
    - Chapter 5 must be a summary of the entire module

    Respond in JSON format with this structure:
    {{
        "chapters": [
            {{
                "title": "Chapter 1 Title",
                "is_summary": false
            }},
            {{
                "title": "Chapter 2 Title",
                "is_summary": false
            }},
            {{
                "title": "Chapter 3 Title",
                "is_summary": false
            }},
            {{
                "title": "Chapter 4 Title",
                "is_summary": false
            }},
            {{
                "title": "Module Summary",
                "is_summary": true
            }}
        ]
    }}
    
    Ensure that each chapter builds on knowledge from previous chapters, creating a logical learning progression.
    """
    
    full_prompt = f"{system_prompt}\n\n{context}"
    
    try:
        logger.info("Sending request to primary HF model (DeepSeek-R1)")
        logger.debug(f"Request model: deepseek-ai/DeepSeek-R1")
        logger.debug(f"Request provider: together")
        
        completion = client.chat.completions.create(
            model="deepseek-ai/DeepSeek-R1",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": context}
            ],
            max_tokens=1000
        )
        response_text = completion.choices[0].message.content
        logger.info(f"Received response from primary model, length: {len(response_text)} characters")
        logger.debug(f"Response preview: {response_text[:100]}...")
    except Exception as e:
        logger.warning(f"Primary model failed, using fallback: {e}")
        try:
            logger.info("Sending request to fallback HF model (Qwen)")
            logger.debug(f"Request model: Qwen/Qwen2-72B-Instruct")
            logger.debug(f"Request provider: hf-inference")
            
            completion = fallback_client.chat.completions.create(
                model="Qwen/Qwen2-72B-Instruct",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": context}
                ],
                max_tokens=1000
            )
            response_text = completion.choices[0].message.content
            logger.info(f"Received response from fallback model, length: {len(response_text)} characters")
            logger.debug(f"Response preview: {response_text[:100]}...")
        except Exception as e2:
            logger.error(f"Both models failed: {e2}")
            return {"error": "Failed to generate chapters due to API issues."}
    
    # Try to extract JSON from the response
    try:
        # First attempt - direct parsing
        result = json.loads(response_text)
    except json.JSONDecodeError:
        # Second attempt - search for JSON in the text
        try:
            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            if json_start >= 0 and json_end > json_start:
                json_text = response_text[json_start:json_end]
                result = json.loads(json_text)
            else:
                raise ValueError("No valid JSON found in the response")
        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"Failed to parse JSON from response: {e}")
            return {"error": "Failed to parse the generated chapters."}
    
    return result

def generate_chapter_content(course_title, module_title, chapter_title, is_summary=False, context_documents=None):
    """
    Generate detailed content for a chapter, including quiz questions if not a summary chapter
    """
    logger.info(f"Generating content for chapter: {chapter_title} in module: {module_title}")
    client = get_hf_inference_client()
    fallback_client = get_hf_fallback_client()
    
    # Create a context based on documents if provided
    context = ""
    if context_documents:
        context = "Based on the following context documents:\n\n"
        for doc in context_documents:
            context += f"{doc.content_text}\n\n"
        logger.info(f"Using {len(context_documents)} context documents for content generation")
    else:
        logger.info("No context documents provided")
    
    if is_summary:
        logger.info("Generating summary content")
        system_prompt = f"""You are an expert cybersecurity educator creating content for a course titled "{course_title}".
        Create a comprehensive summary of the module "{module_title}" for the chapter "{chapter_title}".
        
        This summary should:
        1. Recap key points and concepts from the entire module
        2. Highlight the most important takeaways
        3. Connect concepts across chapters
        4. Prepare learners for applying their knowledge
        
        Format your response in HTML for direct inclusion in the course platform. Use appropriate headings, 
        paragraphs, bullet points, and emphasis as needed. No quiz questions are needed for this summary.
        
        Respond with HTML content only, no JSON formatting necessary.
        """
    else:
        logger.info("Generating regular chapter content with quiz questions")
        system_prompt = f"""You are an expert cybersecurity educator creating content for a course titled "{course_title}".
        Create detailed, comprehensive content for the chapter "{chapter_title}" in the module "{module_title}".
        
        Your content should:
        1. Be thorough and educational, covering all relevant aspects of the topic
        2. Include practical examples and real-world applications in detail and also use at least 500 words 

        3. Be formatted with clear sections, subsections, and visual organization
        
        Additionally, create exactly 3 multiple-choice quiz questions to test understanding of this chapter.
        Each question should have 4 possible answers with exactly one correct answer.
        
        Respond in JSON format with this structure:
        {{
            "chapter_content": "Detailed HTML content for the chapter including formatting",
            "questions": [
                {{
                    "question_text": "Question 1 text",
                    "choices": [
                        {{"choice_text": "Option A", "is_correct": false}},
                        {{"choice_text": "Option B", "is_correct": true}},
                        {{"choice_text": "Option C", "is_correct": false}},
                        {{"choice_text": "Option D", "is_correct": false}}
                    ]
                }},
                // Questions 2 and 3 follow the same structure
            ]
        }}
        
        For the chapter_content, use appropriate HTML tags for formatting (headings, paragraphs, lists, etc.).
        """
    
    full_prompt = f"{system_prompt}\n\n{context}"
    logger.info(f"Prompt length: {len(full_prompt)} characters")
    
    try:
        logger.info("Sending request to primary HF model (DeepSeek-R1)")
        logger.debug(f"Request model: deepseek-ai/DeepSeek-R1")
        logger.debug(f"Request provider: together")
        
        completion = client.chat.completions.create(
            model="deepseek-ai/DeepSeek-R1",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": context}
            ],
            max_tokens=2000
        )
        response_text = completion.choices[0].message.content
        logger.info(f"Received response from primary model, length: {len(response_text)} characters")
        logger.debug(f"Response preview: {response_text[:100]}...")
    except Exception as e:
        logger.warning(f"Primary model failed with error: {str(e)}")
        try:
            logger.info("Sending request to fallback HF model (Qwen)")
            logger.debug(f"Request model: Qwen/Qwen2-72B-Instruct")
            logger.debug(f"Request provider: hf-inference")
            
            completion = fallback_client.chat.completions.create(
                model="Qwen/Qwen2-72B-Instruct",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": context}
                ],
                max_tokens=2000
            )
            response_text = completion.choices[0].message.content
            logger.info(f"Received response from fallback model, length: {len(response_text)} characters")
            logger.debug(f"Response preview: {response_text[:100]}...")
        except Exception as e2:
            logger.error(f"Both models failed. Fallback error: {str(e2)}")
            return {"error": f"Failed to generate chapter content due to API issues. Primary error: {str(e)}, Fallback error: {str(e2)}"}
    
    if is_summary:
        # For summary, return the HTML content directly
        logger.info(f"Successfully generated summary content of length {len(response_text)}")
        return {"chapter_content": response_text, "is_summary": True}
    else:
        # For regular chapters, parse the JSON response
        try:
            # First attempt - direct parsing
            result = json.loads(response_text)
            logger.info(f"Successfully parsed JSON response for chapter content")
            
            # Validate that each question has exactly one correct answer
            if "questions" in result:
                for i, question in enumerate(result["questions"]):
                    correct_count = sum(1 for choice in question["choices"] if choice.get("is_correct") is True)
                    
                    if correct_count != 1:
                        logger.warning(f"Question {i+1} has {correct_count} correct answers. Fixing...")
                        
                        # Fix: If no correct answer, mark the first one as correct
                        if correct_count == 0 and question["choices"]:
                            question["choices"][0]["is_correct"] = True
                        
                        # Fix: If multiple correct answers, keep only the first one
                        elif correct_count > 1:
                            found_correct = False
                            for choice in question["choices"]:
                                if choice.get("is_correct") is True:
                                    if found_correct:
                                        choice["is_correct"] = False
                                    else:
                                        found_correct = True
            
        except json.JSONDecodeError:
            # Second attempt - search for JSON in the text
            try:
                json_start = response_text.find('{')
                json_end = response_text.rfind('}') + 1
                if json_start >= 0 and json_end > json_start:
                    json_text = response_text[json_start:json_end]
                    result = json.loads(json_text)
                    logger.info(f"Extracted JSON from text response for chapter content")
                else:
                    # If we can't find valid JSON, just return the content as HTML
                    logger.warning(f"No valid JSON found, returning raw content as HTML")
                    return {"chapter_content": response_text, "questions": []}
            except (json.JSONDecodeError, ValueError) as e:
                logger.error(f"Failed to parse JSON from response: {e}")
                # Return the raw text as chapter content as a fallback
                return {"chapter_content": response_text, "questions": []}
        
        # Ensure the result has the expected structure
        if "chapter_content" not in result:
            logger.warning("Response doesn't contain chapter_content, adding it")
            result["chapter_content"] = "Content not properly generated. Please try again."
        
        return result

def extract_text_from_document(document_file):
    """Extract text content from uploaded documents (PDF, DOCX, TXT, PPT, PPTX)"""
    filename = document_file.name.lower()
    content = ""
    
    try:
        # Reset file pointer to beginning to ensure we can read the full content
        if hasattr(document_file, 'seek'):
            document_file.seek(0)
            
        if filename.endswith('.pdf'):
            try:
                pdf_file = PyPDF2.PdfReader(document_file)
                for page_num in range(len(pdf_file.pages)):
                    page_text = pdf_file.pages[page_num].extract_text() or ""
                    content += page_text + "\n\n"
                    
                # If we couldn't extract text (might be scanned PDF)
                if not content.strip():
                    return "PDF appears to contain scanned content that couldn't be extracted. Please provide a text-based PDF."
            except Exception as pdf_error:
                logger.error(f"PDF processing error: {pdf_error}")
                return f"Error processing PDF file: {str(pdf_error)}"
        
        elif filename.endswith('.docx'):
            try:
                doc = docx.Document(document_file)
                
                # Extract text from paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():  # Only add non-empty paragraphs
                        content += para.text + "\n"
                
                # Extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            content += row_text + "\n"
                            
            except Exception as docx_error:
                logger.error(f"DOCX processing error: {docx_error}")
                return f"Error processing DOCX file: {str(docx_error)}"
        
        elif filename.endswith('.txt'):
            try:
                # For text files, decode with handling for potential encoding issues
                try:
                    content = document_file.read().decode('utf-8')
                except UnicodeDecodeError:
                    # Try with a different encoding if utf-8 fails
                    document_file.seek(0)
                    content = document_file.read().decode('latin-1')
            except Exception as txt_error:
                logger.error(f"TXT processing error: {txt_error}")
                return f"Error processing TXT file: {str(txt_error)}"
        
        elif (filename.endswith('.pptx') or filename.endswith('.ppt')) and PPTX_AVAILABLE:
            try:
                # For PPT files we need to extract text from slides
                presentation = Presentation(document_file)
                
                # Extract text from each slide
                for slide_number, slide in enumerate(presentation.slides, 1):
                    content += f"Slide {slide_number}:\n"
                    
                    # Extract text from shapes (includes text boxes, titles, etc.)
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            content += f"{shape.text.strip()}\n"
                    
                    # Add extra newline between slides
                    content += "\n\n"
                
                if not content.strip():
                    return "PowerPoint file appears to contain no extractable text. Please check if it contains text content."
                    
            except Exception as ppt_error:
                logger.error(f"PowerPoint processing error: {ppt_error}")
                return f"Error processing PowerPoint file: {str(ppt_error)}"
        elif filename.endswith('.pptx') or filename.endswith('.ppt'):
            return "PowerPoint files are supported but the required library (python-pptx) is not installed."
        else:
            return "Unsupported file format. Please upload PDF, DOCX, TXT, PPT, or PPTX files."
        
        # Check if we actually got content
        if not content.strip():
            return "Could not extract text content from the document. Please check the file format and contents."
            
        # Log size of extracted content
        logger.info(f"Extracted {len(content)} characters from {filename}")
        return content
        
    except Exception as e:
        logger.error(f"Error extracting text from document: {e}")
        return f"Error processing document: {str(e)}"
